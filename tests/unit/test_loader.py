import io

import pymupdf
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.ingestion import loader
from app.ingestion.loader import LOADERS, load_document, load_pdf, load_pptx, load_text


def _tiny_png_bytes() -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", (4, 4), color="white").save(buffer, format="PNG")
    return buffer.getvalue()


# --- load_pdf -----------------------------------------------------------------


def test_load_pdf_extracts_real_text_without_calling_ocr(tmp_path, monkeypatch):
    called = []
    monkeypatch.setattr(loader, "_ocr_image_bytes", lambda b: called.append(b) or "should not be used")

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Employees receive 20 days of annual leave per year.")
    pdf_path = tmp_path / "handbook.pdf"
    doc.save(str(pdf_path))
    doc.close()

    pages = load_pdf(str(pdf_path))

    assert len(pages) == 1
    assert "20 days of annual leave" in pages[0]["text"]
    assert pages[0]["page"] == 1
    assert called == []  # real text present - OCR fallback never triggered


def test_load_pdf_ocrs_a_scanned_page_with_no_text_layer(tmp_path, monkeypatch):
    monkeypatch.setattr(loader, "_ocr_image_bytes", lambda image_bytes: "OCR'd scanned text")

    doc = pymupdf.open()
    doc.new_page()  # blank page, no text layer at all
    pdf_path = tmp_path / "scanned.pdf"
    doc.save(str(pdf_path))
    doc.close()

    pages = load_pdf(str(pdf_path))

    assert pages[0]["text"] == "OCR'd scanned text"


def test_load_pdf_scanned_page_does_not_double_ocr_its_own_image(tmp_path, monkeypatch):
    """A scanned page *is* an embedded image with no separate text layer - the
    page-level OCR fallback and the per-embedded-image OCR pass would
    otherwise both fire on the exact same pixels, duplicating the text."""
    calls = []
    monkeypatch.setattr(loader, "_ocr_image_bytes", lambda b: calls.append(b) or "Sabbatical leave policy")

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_image(pymupdf.Rect(0, 0, 200, 100), stream=_tiny_png_bytes())
    pdf_path = tmp_path / "scanned_with_image.pdf"
    doc.save(str(pdf_path))
    doc.close()

    pages = load_pdf(str(pdf_path))

    assert len(calls) == 1  # OCR'd once (the page render), not again per embedded image
    assert pages[0]["text"].count("Sabbatical leave policy") == 1
    assert "[Image 1]" not in pages[0]["text"]


def test_load_pdf_appends_table_text(tmp_path, monkeypatch):
    monkeypatch.setattr(loader, "_ocr_image_bytes", lambda b: "")

    class FakeTable:
        def extract(self):
            return [["Metric", "Value"], ["Faithfulness", "0.9"]]

    class FakeTableFinder:
        tables = [FakeTable()]

    monkeypatch.setattr(pymupdf.Page, "find_tables", lambda self: FakeTableFinder())

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Results summary")
    pdf_path = tmp_path / "report.pdf"
    doc.save(str(pdf_path))
    doc.close()

    pages = load_pdf(str(pdf_path))

    assert "[Table]" in pages[0]["text"]
    assert "Faithfulness\t0.9" in pages[0]["text"]
    assert "Results summary" in pages[0]["text"]


def test_load_pdf_ocrs_embedded_images(tmp_path, monkeypatch):
    monkeypatch.setattr(loader, "_ocr_image_bytes", lambda image_bytes: "text found in the embedded image")

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Chart below shows quarterly growth.")
    page.insert_image(pymupdf.Rect(72, 100, 200, 200), stream=_tiny_png_bytes())
    pdf_path = tmp_path / "with_chart.pdf"
    doc.save(str(pdf_path))
    doc.close()

    pages = load_pdf(str(pdf_path))

    assert "[Image 1]" in pages[0]["text"]
    assert "text found in the embedded image" in pages[0]["text"]
    assert "quarterly growth" in pages[0]["text"]


def test_load_pdf_table_extraction_failure_does_not_crash_ingestion(tmp_path, monkeypatch):
    monkeypatch.setattr(loader, "_ocr_image_bytes", lambda b: "")
    monkeypatch.setattr(pymupdf.Page, "find_tables", lambda self: (_ for _ in ()).throw(RuntimeError("boom")))

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Still readable despite the table extractor failing.")
    pdf_path = tmp_path / "flaky_table.pdf"
    doc.save(str(pdf_path))
    doc.close()

    pages = load_pdf(str(pdf_path))

    assert "Still readable" in pages[0]["text"]


# --- load_pptx ------------------------------------------------------------------


def test_load_pptx_extracts_text_table_and_image(tmp_path, monkeypatch):
    monkeypatch.setattr(loader, "_ocr_image_bytes", lambda image_bytes: "slide image OCR text")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank layout

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text_frame.text = "Q3 Retrieval Quality Review"

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "Metric"
    table_shape.table.cell(0, 1).text = "Score"
    table_shape.table.cell(1, 0).text = "Recall"
    table_shape.table.cell(1, 1).text = "0.9"

    image_path = tmp_path / "chart.png"
    image_path.write_bytes(_tiny_png_bytes())
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(3), Inches(2), Inches(2))

    pptx_path = tmp_path / "deck.pptx"
    presentation.save(str(pptx_path))

    pages = load_pptx(str(pptx_path))

    assert len(pages) == 1
    text = pages[0]["text"]
    assert "Q3 Retrieval Quality Review" in text
    assert "[Table]" in text
    assert "Recall\t0.9" in text
    assert "[Image]" in text
    assert "slide image OCR text" in text


def test_load_pptx_includes_speaker_notes(tmp_path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = "Remember to mention the reranker."

    pptx_path = tmp_path / "notes.pptx"
    presentation.save(str(pptx_path))

    pages = load_pptx(str(pptx_path))

    assert "[Speaker notes]" in pages[0]["text"]
    assert "reranker" in pages[0]["text"]


# --- load_text ------------------------------------------------------------------


def test_load_text_reads_plain_text_as_a_single_page(tmp_path):
    txt_path = tmp_path / "policy.txt"
    txt_path.write_text("Remote work is permitted up to two days per week.", encoding="utf-8")

    pages = load_text(str(txt_path))

    assert pages == [{"page": 1, "text": "Remote work is permitted up to two days per week."}]


# --- load_document dispatcher ----------------------------------------------------


def test_load_document_dispatches_by_extension(tmp_path):
    txt_path = tmp_path / "notes.txt"
    txt_path.write_text("hello", encoding="utf-8")

    assert load_document(str(txt_path), "notes.txt") == [{"page": 1, "text": "hello"}]


def test_load_document_is_case_insensitive_on_extension(tmp_path):
    txt_path = tmp_path / "notes.txt"
    txt_path.write_text("hello", encoding="utf-8")

    assert load_document(str(txt_path), "NOTES.TXT") == [{"page": 1, "text": "hello"}]


def test_load_document_rejects_unsupported_extension(tmp_path):
    with pytest.raises(ValueError, match="unsupported file type"):
        load_document(str(tmp_path / "whatever.docx"), "whatever.docx")


def test_loaders_registry_covers_pdf_pptx_txt():
    assert set(LOADERS) == {".pdf", ".pptx", ".txt"}
