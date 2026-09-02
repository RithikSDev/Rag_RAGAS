import io
import logging
from pathlib import Path

import fitz
import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger("app.ingestion.loader")

# A page/slide with less real extracted text than this is treated as a
# scanned/image-only page and gets OCR'd instead of indexed empty.
MIN_TEXT_CHARS_BEFORE_OCR = 20


def _ocr_image_bytes(image_bytes: bytes) -> str:
    """OCRs a single embedded image. Isolated in its own function (rather than
    inlined) so tests can monkeypatch it instead of needing a real tesseract
    binary. Never raises - a bad/corrupt embedded image just contributes no
    text, it shouldn't fail the whole document."""
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            return pytesseract.image_to_string(image).strip()
    except Exception:
        logger.warning("OCR failed on an embedded image, skipping it", exc_info=True)
        return ""


def _table_to_text(rows: list[list[str | None]]) -> str:
    return "\n".join("\t".join(cell or "" for cell in row) for row in rows)


def load_pdf(file_path: str) -> list[dict]:
    document = fitz.open(file_path)
    pages = []

    for page_number, page in enumerate(document):
        text = page.get_text().strip()

        # Scanned pages have no extractable text layer - render the page to
        # an image and OCR it instead of indexing it empty. A page-level OCR
        # pass already covers every embedded image on the page (they're baked
        # into the rendered pixmap), so it also stands in for the per-image
        # OCR pass below rather than running OCR on the same content twice.
        page_was_ocrd = False

        if len(text) < MIN_TEXT_CHARS_BEFORE_OCR:
            ocr_text = _ocr_image_bytes(page.get_pixmap(dpi=200).tobytes("png"))
            if ocr_text:
                text = f"{text}\n{ocr_text}".strip()
                page_was_ocrd = True

        # Tables carry structure plain text extraction loses - append them as
        # tab-separated rows so retrieval can still match on cell values.
        try:
            for table in page.find_tables().tables:
                table_text = _table_to_text(table.extract())
                if table_text.strip():
                    text = f"{text}\n\n[Table]\n{table_text}".strip()
        except Exception:
            logger.warning("table extraction failed on page %s", page_number + 1, exc_info=True)

        # Embedded images (screenshots, charts, photographed text) can carry
        # text a page-text pass never sees - OCR each one and fold it in.
        if not page_was_ocrd:
            for image_index, image_info in enumerate(page.get_images(full=True)):
                try:
                    image_bytes = document.extract_image(image_info[0])["image"]
                    ocr_text = _ocr_image_bytes(image_bytes)
                    if ocr_text:
                        text = f"{text}\n\n[Image {image_index + 1}]\n{ocr_text}".strip()
                except Exception:
                    logger.warning(
                        "failed to OCR embedded image %s on page %s",
                        image_index + 1,
                        page_number + 1,
                        exc_info=True,
                    )

        pages.append({"page": page_number + 1, "text": text})

    document.close()

    return pages


def load_pptx(file_path: str) -> list[dict]:
    presentation = Presentation(file_path)
    pages = []

    for slide_number, slide in enumerate(presentation.slides):
        parts = []

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())

            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                table_text = _table_to_text(rows)
                if table_text.strip():
                    parts.append(f"[Table]\n{table_text}")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ocr_text = _ocr_image_bytes(shape.image.blob)
                if ocr_text:
                    parts.append(f"[Image]\n{ocr_text}")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes]\n{notes}")

        pages.append({"page": slide_number + 1, "text": "\n\n".join(parts)})

    return pages


def load_text(file_path: str) -> list[dict]:
    text = Path(file_path).read_text(encoding="utf-8")
    return [{"page": 1, "text": text}]


LOADERS = {
    ".pdf": load_pdf,
    ".pptx": load_pptx,
    ".txt": load_text,
}


def load_document(file_path: str, original_filename: str) -> list[dict]:
    suffix = Path(original_filename).suffix.lower()
    loader = LOADERS.get(suffix)

    if loader is None:
        raise ValueError(f"unsupported file type: {suffix}")

    return loader(file_path)
